import streamlit as st
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import os

st.title("🤖 ChatBot de Bioquímica – Aminoácidos")
st.success("👩‍⚕️ Bienvenido/a al ChatBot de Bioquímica.\n\nEste asistente responde preguntas sobre aminoácidos usando:\n- Las diapositivas de clase\n- Un capítulo del libro\n- El video de la Dra. Susana González Chávez")

# Archivos fuente
pptx_path = "clase_001_aminoacidos.pptx"
txt_path = "capitulo_aminoacidos_mckee.txt"

# Video relacionado
videos = {
    "clase_001_aminoacidos.pptx": "https://youtu.be/6-rvZqSTANo?si=WfT34ODacliTwOhz"
}

# Función para extraer texto del pptx
def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    slides_text = []
    for slide in prs.slides:
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
        slides_text.append(text.strip())
    return slides_text

# Función para extraer texto del txt (capítulo del libro)
def extract_text_from_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        content = f.read()
    # Dividir en párrafos para mejor análisis
    paragraphs = [p.strip() for p in content.split("\n") if len(p.strip()) > 50]
    return paragraphs

# Cargar contenido
try:
    slides = extract_text_from_pptx(pptx_path)
except Exception as e:
    st.error(f"Error al leer el archivo de diapositivas: {e}")
    st.stop()

try:
    chapter = extract_text_from_txt(txt_path)
except Exception as e:
    st.error(f"Error al leer el archivo del capítulo: {e}")
    st.stop()

# Entrada del usuario
query = st.text_input("Escribe tu pregunta sobre aminoácidos:")

if query:
    # Combinar fuentes
    all_docs = slides + chapter
    vectorizer = TfidfVectorizer().fit_transform([query] + all_docs)
    similarity = cosine_similarity(vectorizer[0:1], vectorizer[1:])
    best_idx = similarity.argmax()

    st.subheader("📖 Respuesta basada en tus materiales:")
    st.write(all_docs[best_idx])

    # Mostrar video si existe
    video_url = videos.get(pptx_path)
    if video_url:
        st.markdown("🎥 **También puedes ver la explicación en video:**")
        st.video(video_url)
