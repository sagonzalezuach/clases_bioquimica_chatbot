import streamlit as st
import os
from openai import OpenAI
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

# Inicializar cliente OpenAI
client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

# --- Personalización visual ---
st.set_page_config(page_title="ChatBot Bioquímica UACH", page_icon="🔖", layout="centered")

# Buscar minuto relevante en temas_video
tema_encontrado = None
query_lower = query.lower()

for tema, minuto in temas_video.items():
    if tema in query_lower:
        tema_encontrado = (tema, minuto)
        break

if tema_encontrado:
    st.markdown(f"🎯 **Puedes encontrar la explicación de *{tema_encontrado[0]}* en el minuto {tema_encontrado[1]} del video.**")
st.markdown("""
    <style>
    body {
        background-color: #fff0f5;
    }
    .main {
        background-color: #ffe6f0;
        padding: 2rem;
        border-radius: 12px;
    }
    h1, h2, h3, h4 {
        color: #c2185b;
    }
    .stButton button {
        background-color: #ec407a;
        color: white;
    }
    </style>
""", unsafe_allow_html=True)

# --- Encabezado con logo local ---
st.image("logo uach.png", width=100)

# --- Banner decorativo ---
st.markdown("""
<div style="background-color:#ffc0cb;padding:20px;border-radius:10px;text-align:center">
    <h1 style="color:#880e4f;">ChatBot de Bioquímica</h1>
    <h3 style="color:#6a1b9a;">Facultad de Medicina y Ciencias Biomédicas</h3>
    <h4 style="color:#6a1b9a;">Universidad Autónoma de Chihuahua</h4>
    <p style="color:#4a148c;">Este asistente responde preguntas sobre aminoácidos usando tus propias clases: presentaciones, lectura y video.</p>
</div>
""", unsafe_allow_html=True)

# Archivos fuente (ajustados para estar en raíz del repositorio)
pptx_path = "clase_001_aminoacidos.pptx"
txt_path = "capitulo_aminoacidos_mckee_LIMPIO.txt"
video_url = "https://youtu.be/6-rvZqSTANo?si=WfT34ODacliTwOhz"
# Tabla de temas y minutos del video
temas_video = {
    "estructura general": "8:33",
    "tipos de aminoácidos": "11:38",
    "aminoácidos polares": "31:11",
    "aminoácidos apolares": "21:07",
    "aminoácidos ácidos": "35:33"
}


# Funciones para cargar contenido
def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    return [" ".join(shape.text for shape in slide.shapes if hasattr(shape, "text")).strip() for slide in prs.slides]

def extract_text_from_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return [p.strip() for p in f.read().split("\n\n") if len(p.strip()) > 60]

slides = extract_text_from_pptx(pptx_path)
chapter = extract_text_from_txt(txt_path)
all_docs = slides + chapter

# Entrada del usuario
query = st.text_input("🎓 Escribe tu pregunta sobre aminoácidos:")

if query:
    # Sugerencias por minuto
    temas_video = {
        "estructura general": "8:33",
        "tipos de aminoácidos": "11:38",
        "aminoácidos polares": "31:11",
        "aminoácidos apolares": "21:07",
        "aminoácidos ácidos": "35:33"
    }

    query_lower = query.lower()
    tema_encontrado = None
    for tema, minuto in temas_video.items():
        if tema in query_lower:
            tema_encontrado = (tema, minuto)
            break

    # Vectorización y búsqueda de contexto
    vectorizer = TfidfVectorizer().fit_transform([query] + all_docs)
    similarity = cosine_similarity(vectorizer[0:1], vectorizer[1:])
    top_indices = similarity[0].argsort()[-3:][::-1]
    context = "\n\n".join([all_docs[i] for i in top_indices])

    # Construir prompt
    prompt = f"""
Eres un asistente de bioquímica de la Facultad de Medicina de la UACH.
Responde solo usando los siguientes materiales proporcionados por la Dra. Susana González Chávez.
No inventes nada fuera del contenido.

PREGUNTA: {query}

MATERIALES:
{context}

RESPUESTA:
"""

    with st.spinner("🤖 Consultando a GPT-4..."):
        response = client.chat.completions.create(
            model="gpt-4",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.2,
            max_tokens=600
        )

        st.subheader(":open_book: Respuesta elaborada con base en tus clases:")
        st.write(response.choices[0].message.content)
        st.caption("📚 Esta respuesta se generó con base en tus presentaciones, lectura y video. No se utilizó información externa.")

        if tema_encontrado:
            st.markdown(f"🎯 **Puedes encontrar la explicación de *{tema_encontrado[0]}* en el minuto {tema_encontrado[1]} del video.**")

        st.markdown("**🎥 Explicación en video:**")
        st.video(video_url)


# Pie de página
st.markdown("""
---
**Desarrollado por la Dra. Susana González Chávez**  
Facultad de Medicina y Ciencias Biomédicas, UACH  
:heart: Proyecto educativo con fines docentes
""")
